"""PowerPoint (.pptx) to markdown reader."""

from __future__ import annotations

import sys
from pathlib import Path


def to_markdown(
    path: Path,
    *,
    head: int | None = None,
    slide: int | None = None,
    **_kwargs,
) -> str:
    """Convert a pptx file to markdown."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.exc import PackageNotFoundError

    try:
        prs = Presentation(str(path))
    except PackageNotFoundError:
        from officecat.readers._errors import report_open_error

        report_open_error(path)
        raise SystemExit(3)
    except Exception as e:
        msg = str(e).lower()
        if "password" in msg or "encrypt" in msg:
            print(
                f"Error: '{path.name}' is password-protected. "
                f"officecat cannot open encrypted files.",
                file=sys.stderr,
            )
            raise SystemExit(3)
        raise

    total_slides = len(prs.slides)

    if slide is not None:
        if slide < 1 or slide > total_slides:
            print(
                f"Error: Slide {slide} not found. "
                f"Document has {total_slides} slides.",
                file=sys.stderr,
            )
            raise SystemExit(1)

    sections: list[str] = []

    for i, sld in enumerate(prs.slides, start=1):
        if slide is not None and i != slide:
            continue

        # Check if slide is hidden
        show_attr = sld._element.get("show")
        is_hidden = show_attr == "0"

        title_shape = sld.shapes.title
        title_text = title_shape.text.strip() if title_shape else None
        title_shape_id = title_shape.shape_id if title_shape else None

        hidden_tag = " (Hidden)" if is_hidden else ""
        if title_text:
            heading = f"## Slide {i}{hidden_tag}: {title_text}"
        else:
            heading = f"## Slide {i}{hidden_tag}"

        body_lines: list[str] = [heading, ""]

        for shape in sld.shapes:
            if title_shape_id is not None and shape.shape_id == title_shape_id:
                continue

            st = shape.shape_type
            if st == MSO_SHAPE_TYPE.PICTURE:
                body_lines.append(f"*[Image: {shape.name}]*")
            elif st == MSO_SHAPE_TYPE.GROUP:
                body_lines.append("*[Grouped content]*")
            elif hasattr(shape, "has_table") and shape.has_table:
                tbl = shape.table
                rows = len(tbl.rows)
                cols = len(tbl.columns)
                body_lines.append(f"*[Table: {rows}x{cols}]*")
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    body_lines.append(text)

        # Notes
        if sld.has_notes_slide:
            notes_text = sld.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                body_lines.append("")
                body_lines.append(f"> **Notes:** {notes_text}")

        sections.append("\n".join(body_lines))

        if slide is not None:
            break
        if head is not None and len(sections) >= head:
            break

    return "\n\n---\n\n".join(sections)
